
import collections.abc
from pptx import Presentation
import camelot
import io
import os


def pdfTableExtract(pdf_file):
    tables = camelot.read_pdf(pdf_file)

    # number of tables extracted
    print("Total tables extracted:", tables.n)

    if tables.n:
        for i in range(tables.n):
            table_df = tables[i].df
            print(table_df)


def writeFile(text):
    for key, value in text.items():
        with io.open(key + ".txt", 'w+', encoding='utf8') as f:
            f.write(value)


def powerPointExtract(ppPath):
    head_tail = os.path.split(ppPath)
    assert head_tail[1].split('.')[1] == 'pptx', 'The file extension is not pptx'

    ppt = Presentation(ppPath)
    text_runs = {}

    for page, slide in enumerate(ppt.slides):
        textNote = slide.notes_slide.notes_text_frame.text
        if textNote:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "PPT" in run.text:
                            file_name = run.text.split(": ")[1]
                            text_runs[file_name] = textNote
    writeFile(text_runs)

if __name__ == '__main__':
    pp_file = "HKSIP1E.pptx"
    pdf_file = "data.pdf"

    pdfTableExtract(pdf_file)

    powerPointExtract(pp_file)
